#!/usr/bin/env python3
"""
create_topic3_pptx.py

Generates Topic3.pptx (20 slides) with a simple/light PowerPoint theme and speaker notes.
Requires: python-pptx
Install: pip install python-pptx
Run: python create_topic3_pptx.py
Output: Topic3.pptx in the script directory
"""

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()  # default (simple/light) template
title_layout = prs.slide_layouts[0]      # Title Slide layout
content_layout = prs.slide_layouts[1]    # Title and Content layout

def add_title_slide(prs, title, subtitle=None, notes_text=""):
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if subtitle:
        try:
            slide.placeholders[1].text = subtitle
        except Exception:
            pass
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

def add_bullet_slide(prs, title, bullets, notes_text=""):
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text

slides = [
    ("Exploring IT Careers — Industries, Roles, and Personal Development",
     ["Course Level: First university course for IT specialists",
      "CLO1: Analyze industries where IT skills are in demand and impact sectors",
      "CLO2: Evaluate strengths and create a professional development plan",
      "Activity: Compare Your Occupational Options"],
     "Introduce course aim and CLOs. Set expectations: industry analysis + self-assessment + actionable plan."),
    ("Learning Objectives & Session Map",
     ["Identify top industries hiring IT specialists and why (CLO1)",
      "Match IT roles to industry needs (CLO1)",
      "Conduct a personal strengths assessment and map to career paths (CLO2)",
      "Draft a short-term professional development plan (CLO2)"],
     "Explain the learning path and assessment of outcomes (participation, activity worksheet, short plan)."),
    ("Why Industry Analysis Matters for IT Specialists",
     ["Industries differ in problems, tech stacks, and compliance needs",
      "Benefits: better job fit, targeted skill development, stronger applications"],
     "Give quick example: cybersecurity priorities in finance vs. healthcare compliance constraints."),
    ("Major Industry Categories with High IT Demand",
     ["Finance & FinTech, Healthcare & HealthTech, Technology & Software Services",
      "Manufacturing, Retail & eCommerce, Telecom, Energy, Government, Education, Media/Gaming"],
     "Briefly describe why each industry values IT skills: data, automation, security, customer experience."),
    ("Industry Impact Examples (Finance & Healthcare)",
     ["Finance: trading, fraud detection, secure transactions (cloud, ML, security)",
      "Healthcare: EHRs, telemedicine, devices (interop, privacy, embedded)"],
     "Use a case: hospital adopting telehealth needs HIPAA-aware devs and secure cloud setups."),
    ("Industry Impact Examples (Manufacturing & Retail)",
     ["Manufacturing: IoT, predictive maintenance, automation (embedded, analytics)",
      "Retail: personalized commerce, inventory automation, omni-channel UX & APIs"],
     "Relate to student experience: online shopping personalization vs. factory sensors."),
    ("In-demand IT Roles Across Industries",
     ["Software Developer, Data Scientist/Engineer, Cloud/DevOps, Cybersecurity Analyst",
      "Sys/Net Admin, Product Manager, UX/UI Designer, Embedded Engineer, SRE"],
     "Point out role name variations; responsibilities often overlap."),
    ("Skills, Tools & Certifications by Role",
     ["Software: Python/Java/JS, Git, CI/CD",
      "Data: SQL, Spark, ML frameworks; Cloud/DevOps: AWS/Azure/GCP, Kubernetes",
      "Security: pen-testing, risk mgmt; Embedded: C/C++, RTOS"],
     "Encourage students to note transferable skills (programming, problem-solving, teamwork)."),
    ("Mapping Skills to Industry Needs",
     ["Create a skills×industry matrix: demand level, tools, certifications",
      "Use it to prioritize learning based on desired industry"],
     "Show a quick 3x3 example live if possible (e.g., Cloud: High in Tech/Finance/Retail)."),
    ("Emerging Trends Driving Demand",
     ["AI/ML & Generative AI, Cloud-native & serverless, Cybersecurity/Zero-Trust",
      "Automation/RPA, Edge computing & IoT, Green computing"],
     "Discuss how trends shift hiring priorities and create new hybrid roles."),
    ("How to Research an Industry (practical sources)",
     ["Job boards (LinkedIn, Indeed), industry reports (Gartner, McKinsey)",
      "O*NET, company tech blogs, alumni, informational interviews, meetups"],
     "Walk students through a quick demo of a job posting and extracting skill signals."),
    ("Know Yourself: Strengths, Interests, Values",
     ["Self-audit: technical skills you enjoy, soft skills, work preferences, values",
      "Tools: reflection journal, StrengthsFinder, skills inventory"],
     "Encourage honest reflection; no single 'best' path — fit matters."),
    ("Mapping Personal Profile to Career Paths",
     ["Personas: Problem Solver → Software/ML; Systems Builder → DevOps; Analyst → Data",
      "Match: day-to-day tasks, learning curve, industry fit, certs cost/time"],
     "Ask students which persona fits them most."),
    ("Building a Personalized Professional Development Plan",
     ["Components: career target, skills, learning sources, experience milestones, metrics",
      "Example: 6-month goal — finish cloud fundamentals + build portfolio app"],
     "Show a short sample timeline and emphasize measurable goals."),
    ("Credentials, Portfolios & Job-Search Materials",
     ["Which certs matter where (cloud certs, CISSP, CompTIA)",
      "Portfolio: 2–4 quality projects, README, deployment, tests, CI; tailor per industry"],
     "Give dos and don'ts for portfolios and certifications."),
    ("Networking & Hands-on Experience",
     ["Internships, open-source, hackathons, campus projects; volunteer opportunities",
      "Networking: informational interviews, alumni outreach, meetups, maintain GitHub/LinkedIn"],
     "Stress quality over quantity in networking; prepare short elevator pitch."),
    ("Activity: Compare Your Occupational Options (Instructions)",
     ["Objective: research and compare at least two IT options and reflect on fit",
      "Steps: pick roles, research postings/skills/salaries, fill comparison, reflect & choose"],
     "Explain outputs and grading: completeness (table), analysis quality, reflection clarity."),
    ("Activity Worksheet: Comparison Criteria",
     ["Collect: typical employers, duties, technical skills, soft skills, certs, salary range",
      "Also: industry fit and personal fit rating (1–5) with justification"],
     "Provide an example filled row to illustrate expectations."),
    ("Next Steps, Assessment & Resources",
     ["Submit: comparison + reflection + 6-month plan; assessment rubric provided",
      "Resources: LinkedIn, StackOverflow Jobs, Gartner, Coursera, O*NET, Glassdoor"],
     "Conclude tying back to CLO1 and CLO2. Offer office hours for 1:1 plan review.")
]

# First slide as title slide
first = slides[0]
add_title_slide(prs, first[0], subtitle=None, notes_text=first[2])

# Remaining slides
for title, bullets, notes in slides[1:]:
    add_bullet_slide(prs, title, bullets, notes_text=notes)

output_filename = "Topic3.pptx"
prs.save(output_filename)
print(f"Saved presentation as {output_filename}")